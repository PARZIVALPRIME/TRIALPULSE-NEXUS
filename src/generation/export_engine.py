# src/generation/export_engine.py
"""
TRIALPULSE NEXUS 10X - Export Engine v1.0
Phase 6.3: PDF/Word/PPT Export with Professional Styling

Features:
- Multiple PDF backends (WeasyPrint, pdfkit, reportlab fallbacks)
- Professional DOCX styling with headers, footers, TOC
- Branded PPTX templates with master slides
- Consistent styling across all formats
- Export configuration management
"""

import os
import sys
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple, Union
from dataclasses import dataclass, field
from enum import Enum
import hashlib
import json
import tempfile
import shutil

# Add project root to path
PROJECT_ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from config.settings import DATA_OUTPUTS

# =============================================================================
# CONFIGURATION
# =============================================================================

@dataclass
class StyleConfig:
    """Styling configuration for exports."""
    # Brand colors
    primary_color: str = "#1a365d"      # Dark blue
    secondary_color: str = "#2c5282"    # Medium blue
    accent_color: str = "#38a169"       # Green
    warning_color: str = "#dd6b20"      # Orange
    danger_color: str = "#e53e3e"       # Red
    success_color: str = "#38a169"      # Green
    
    # Typography
    heading_font: str = "Arial"
    body_font: str = "Calibri"
    mono_font: str = "Consolas"
    
    # Sizes (in points for documents)
    title_size: int = 24
    heading1_size: int = 18
    heading2_size: int = 14
    heading3_size: int = 12
    body_size: int = 11
    small_size: int = 9
    
    # Page settings
    page_margin_inches: float = 1.0
    header_height_inches: float = 0.5
    footer_height_inches: float = 0.5
    
    # Logo path (optional)
    logo_path: Optional[str] = None
    company_name: str = "TrialPulse Nexus"


@dataclass
class ExportResult:
    """Result of an export operation."""
    success: bool
    file_path: Optional[str]
    format: str
    file_size: int = 0
    generation_time_ms: float = 0
    backend_used: Optional[str] = None
    warnings: List[str] = field(default_factory=list)
    error: Optional[str] = None
    
    def to_dict(self) -> Dict:
        return {
            'success': self.success,
            'file_path': self.file_path,
            'format': self.format,
            'file_size': self.file_size,
            'generation_time_ms': self.generation_time_ms,
            'backend_used': self.backend_used,
            'warnings': self.warnings,
            'error': self.error
        }


class PDFBackend(Enum):
    """Available PDF generation backends."""
    WEASYPRINT = "weasyprint"
    PDFKIT = "pdfkit"
    REPORTLAB = "reportlab"
    XHTML2PDF = "xhtml2pdf"


# =============================================================================
# PDF EXPORTER
# =============================================================================

class PDFExporter:
    """
    PDF export with multiple backend fallbacks.
    
    Priority order:
    1. WeasyPrint (best quality, requires GTK)
    2. pdfkit (requires wkhtmltopdf)
    3. xhtml2pdf (pure Python, basic)
    4. reportlab (pure Python, programmatic)
    """
    
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
        self.available_backends = self._detect_backends()
        
    def _detect_backends(self) -> List[PDFBackend]:
        """Detect which PDF backends are available."""
        available = []
        
        # Check WeasyPrint
        try:
            import weasyprint
            available.append(PDFBackend.WEASYPRINT)
        except ImportError:
            pass
        except OSError:
            # GTK not available
            pass
            
        # Check pdfkit
        try:
            import pdfkit
            # Also need wkhtmltopdf installed
            try:
                pdfkit.configuration()
                available.append(PDFBackend.PDFKIT)
            except OSError:
                pass
        except ImportError:
            pass
            
        # Check xhtml2pdf
        try:
            import xhtml2pdf
            available.append(PDFBackend.XHTML2PDF)
        except ImportError:
            pass
            
        # Check reportlab
        try:
            import reportlab
            available.append(PDFBackend.REPORTLAB)
        except ImportError:
            pass
            
        return available
    
    def get_css_styles(self) -> str:
        """Generate CSS styles for PDF."""
        return f"""
        @page {{
            size: A4;
            margin: {self.style.page_margin_inches}in;
            @top-center {{
                content: "{self.style.company_name}";
                font-family: {self.style.heading_font};
                font-size: {self.style.small_size}pt;
                color: {self.style.secondary_color};
            }}
            @bottom-center {{
                content: "Page " counter(page) " of " counter(pages);
                font-family: {self.style.body_font};
                font-size: {self.style.small_size}pt;
            }}
            @bottom-right {{
                content: "Generated: " string(date);
                font-family: {self.style.body_font};
                font-size: {self.style.small_size}pt;
            }}
        }}
        
        body {{
            font-family: {self.style.body_font}, sans-serif;
            font-size: {self.style.body_size}pt;
            line-height: 1.5;
            color: #333;
        }}
        
        h1 {{
            font-family: {self.style.heading_font}, sans-serif;
            font-size: {self.style.title_size}pt;
            color: {self.style.primary_color};
            border-bottom: 3px solid {self.style.accent_color};
            padding-bottom: 10px;
            margin-bottom: 20px;
        }}
        
        h2 {{
            font-family: {self.style.heading_font}, sans-serif;
            font-size: {self.style.heading1_size}pt;
            color: {self.style.secondary_color};
            margin-top: 25px;
            margin-bottom: 15px;
        }}
        
        h3 {{
            font-family: {self.style.heading_font}, sans-serif;
            font-size: {self.style.heading2_size}pt;
            color: {self.style.primary_color};
            margin-top: 20px;
            margin-bottom: 10px;
        }}
        
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 15px 0;
            font-size: {self.style.body_size}pt;
        }}
        
        th {{
            background-color: {self.style.primary_color};
            color: white;
            padding: 10px 8px;
            text-align: left;
            font-weight: bold;
        }}
        
        td {{
            padding: 8px;
            border-bottom: 1px solid #ddd;
        }}
        
        tr:nth-child(even) {{
            background-color: #f8f9fa;
        }}
        
        .metric-card {{
            background: linear-gradient(135deg, {self.style.primary_color} 0%, {self.style.secondary_color} 100%);
            color: white;
            padding: 15px;
            border-radius: 8px;
            margin: 10px 0;
        }}
        
        .metric-value {{
            font-size: {self.style.title_size}pt;
            font-weight: bold;
        }}
        
        .metric-label {{
            font-size: {self.style.small_size}pt;
            opacity: 0.9;
        }}
        
        .status-pristine {{ color: {self.style.success_color}; font-weight: bold; }}
        .status-excellent {{ color: #2f855a; }}
        .status-good {{ color: #38a169; }}
        .status-fair {{ color: {self.style.warning_color}; }}
        .status-poor {{ color: #c53030; }}
        .status-critical {{ color: {self.style.danger_color}; font-weight: bold; }}
        
        .priority-critical {{ 
            background-color: {self.style.danger_color}; 
            color: white; 
            padding: 2px 8px; 
            border-radius: 4px;
            font-weight: bold;
        }}
        .priority-high {{ 
            background-color: {self.style.warning_color}; 
            color: white; 
            padding: 2px 8px; 
            border-radius: 4px;
        }}
        .priority-medium {{ 
            background-color: #ecc94b; 
            color: #744210; 
            padding: 2px 8px; 
            border-radius: 4px;
        }}
        .priority-low {{ 
            background-color: #48bb78; 
            color: white; 
            padding: 2px 8px; 
            border-radius: 4px;
        }}
        
        .highlight-box {{
            background-color: #ebf8ff;
            border-left: 4px solid {self.style.accent_color};
            padding: 15px;
            margin: 15px 0;
        }}
        
        .warning-box {{
            background-color: #fffaf0;
            border-left: 4px solid {self.style.warning_color};
            padding: 15px;
            margin: 15px 0;
        }}
        
        .danger-box {{
            background-color: #fff5f5;
            border-left: 4px solid {self.style.danger_color};
            padding: 15px;
            margin: 15px 0;
        }}
        
        .footer {{
            margin-top: 30px;
            padding-top: 15px;
            border-top: 1px solid #ddd;
            font-size: {self.style.small_size}pt;
            color: #666;
        }}
        
        .page-break {{
            page-break-before: always;
        }}
        """
    
    def export(
        self,
        html_content: str,
        output_path: str,
        preferred_backend: Optional[PDFBackend] = None
    ) -> ExportResult:
        """
        Export HTML to PDF using available backend.
        
        Args:
            html_content: HTML string to convert
            output_path: Where to save the PDF
            preferred_backend: Preferred backend (falls back if unavailable)
            
        Returns:
            ExportResult with details
        """
        import time
        start_time = time.time()
        
        # Determine backend to use
        if preferred_backend and preferred_backend in self.available_backends:
            backends_to_try = [preferred_backend] + [b for b in self.available_backends if b != preferred_backend]
        else:
            backends_to_try = self.available_backends
            
        if not backends_to_try:
            return ExportResult(
                success=False,
                file_path=None,
                format='pdf',
                error="No PDF backend available. Install one of: weasyprint, pdfkit (+ wkhtmltopdf), xhtml2pdf, reportlab"
            )
        
        # Inject CSS into HTML
        styled_html = self._inject_styles(html_content)
        
        # Try each backend
        for backend in backends_to_try:
            try:
                if backend == PDFBackend.WEASYPRINT:
                    self._export_weasyprint(styled_html, output_path)
                elif backend == PDFBackend.PDFKIT:
                    self._export_pdfkit(styled_html, output_path)
                elif backend == PDFBackend.XHTML2PDF:
                    self._export_xhtml2pdf(styled_html, output_path)
                elif backend == PDFBackend.REPORTLAB:
                    self._export_reportlab(styled_html, output_path)
                
                # Success
                file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
                elapsed = (time.time() - start_time) * 1000
                
                return ExportResult(
                    success=True,
                    file_path=output_path,
                    format='pdf',
                    file_size=file_size,
                    generation_time_ms=elapsed,
                    backend_used=backend.value
                )
                
            except Exception as e:
                continue  # Try next backend
        
        # All backends failed
        elapsed = (time.time() - start_time) * 1000
        return ExportResult(
            success=False,
            file_path=None,
            format='pdf',
            generation_time_ms=elapsed,
            error=f"All PDF backends failed. Tried: {[b.value for b in backends_to_try]}"
        )
    
    def _inject_styles(self, html_content: str) -> str:
        """Inject CSS styles into HTML."""
        css = self.get_css_styles()
        
        if '<head>' in html_content:
            return html_content.replace(
                '<head>',
                f'<head><style>{css}</style>'
            )
        else:
            return f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <style>{css}</style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
    
    def _export_weasyprint(self, html: str, output_path: str):
        """Export using WeasyPrint."""
        from weasyprint import HTML, CSS
        HTML(string=html).write_pdf(output_path)
    
    def _export_pdfkit(self, html: str, output_path: str):
        """Export using pdfkit."""
        import pdfkit
        options = {
            'page-size': 'A4',
            'margin-top': f'{self.style.page_margin_inches}in',
            'margin-right': f'{self.style.page_margin_inches}in',
            'margin-bottom': f'{self.style.page_margin_inches}in',
            'margin-left': f'{self.style.page_margin_inches}in',
            'encoding': 'UTF-8',
            'enable-local-file-access': None
        }
        pdfkit.from_string(html, output_path, options=options)
    
    def _export_xhtml2pdf(self, html: str, output_path: str):
        """Export using xhtml2pdf."""
        from xhtml2pdf import pisa
        with open(output_path, 'wb') as pdf_file:
            pisa_status = pisa.CreatePDF(html, dest=pdf_file)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
    
    def _export_reportlab(self, html: str, output_path: str):
        """Export using reportlab (basic HTML to PDF)."""
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.lib.units import inch
        import re
        
        doc = SimpleDocTemplate(
            output_path,
            pagesize=A4,
            rightMargin=self.style.page_margin_inches * inch,
            leftMargin=self.style.page_margin_inches * inch,
            topMargin=self.style.page_margin_inches * inch,
            bottomMargin=self.style.page_margin_inches * inch
        )
        
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name='CustomTitle',
            fontSize=self.style.title_size,
            spaceAfter=20,
            textColor=colors.HexColor(self.style.primary_color)
        ))
        
        story = []
        
        # Simple HTML parsing for reportlab
        # Extract text content from common tags
        html_clean = re.sub(r'<style.*?</style>', '', html, flags=re.DOTALL)
        html_clean = re.sub(r'<script.*?</script>', '', html_clean, flags=re.DOTALL)
        
        # Find headings
        for match in re.finditer(r'<h1[^>]*>(.*?)</h1>', html_clean, re.DOTALL):
            story.append(Paragraph(match.group(1).strip(), styles['CustomTitle']))
            story.append(Spacer(1, 12))
        
        for match in re.finditer(r'<h2[^>]*>(.*?)</h2>', html_clean, re.DOTALL):
            story.append(Paragraph(match.group(1).strip(), styles['Heading2']))
            story.append(Spacer(1, 8))
        
        for match in re.finditer(r'<p[^>]*>(.*?)</p>', html_clean, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1))
            story.append(Paragraph(text.strip(), styles['Normal']))
            story.append(Spacer(1, 6))
        
        if not story:
            # Fallback: just extract all text
            text = re.sub(r'<[^>]+>', ' ', html_clean)
            text = re.sub(r'\s+', ' ', text).strip()
            story.append(Paragraph(text[:5000], styles['Normal']))
        
        doc.build(story)


# =============================================================================
# DOCX EXPORTER
# =============================================================================

class DOCXExporter:
    """
    Professional Word document exporter with enhanced styling.
    """
    
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
        
    def export(
        self,
        content: Dict[str, Any],
        output_path: str,
        include_toc: bool = True,
        include_header: bool = True,
        include_footer: bool = True
    ) -> ExportResult:
        """
        Export content to a professionally styled DOCX.
        
        Args:
            content: Dictionary with document content:
                - title: Document title
                - subtitle: Optional subtitle
                - sections: List of section dicts with 'heading', 'content', 'level'
                - tables: List of table dicts with 'title', 'headers', 'rows'
                - metadata: Dict with author, date, etc.
            output_path: Where to save the DOCX
            include_toc: Include table of contents
            include_header: Include header with company name
            include_footer: Include footer with page numbers
            
        Returns:
            ExportResult with details
        """
        import time
        start_time = time.time()
        
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.style import WD_STYLE_TYPE
            from docx.enum.table import WD_TABLE_ALIGNMENT
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
        except ImportError:
            return ExportResult(
                success=False,
                file_path=None,
                format='docx',
                error="python-docx not installed. Run: pip install python-docx"
            )
        
        try:
            doc = Document()
            
            # Set up document properties
            core_properties = doc.core_properties
            core_properties.author = content.get('metadata', {}).get('author', 'TrialPulse Nexus')
            core_properties.title = content.get('title', 'Report')
            core_properties.created = datetime.now()
            
            # Set up styles
            self._setup_styles(doc)
            
            # Add header
            if include_header:
                self._add_header(doc, content.get('title', 'Report'))
            
            # Add footer
            if include_footer:
                self._add_footer(doc)
            
            # Title
            title_para = doc.add_paragraph()
            title_run = title_para.add_run(content.get('title', 'Report'))
            title_run.bold = True
            title_run.font.size = Pt(self.style.title_size)
            title_run.font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Subtitle
            if content.get('subtitle'):
                subtitle_para = doc.add_paragraph()
                subtitle_run = subtitle_para.add_run(content['subtitle'])
                subtitle_run.font.size = Pt(self.style.heading2_size)
                subtitle_run.font.color.rgb = RGBColor.from_string(self.style.secondary_color[1:])
                subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Date
            date_para = doc.add_paragraph()
            date_run = date_para.add_run(
                content.get('metadata', {}).get('date', datetime.now().strftime('%B %d, %Y'))
            )
            date_run.font.size = Pt(self.style.body_size)
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph()  # Spacer
            
            # Table of Contents placeholder
            if include_toc:
                doc.add_paragraph('Table of Contents', style='Heading 1')
                toc_para = doc.add_paragraph()
                toc_para.add_run('[Table of Contents - Update field to generate]')
                toc_para.italic = True
                doc.add_page_break()
            
            # Sections
            for section in content.get('sections', []):
                level = section.get('level', 1)
                heading_style = f'Heading {min(level, 3)}'
                
                doc.add_paragraph(section.get('heading', ''), style=heading_style)
                
                # Content can be string or list of paragraphs
                section_content = section.get('content', '')
                if isinstance(section_content, str):
                    doc.add_paragraph(section_content)
                elif isinstance(section_content, list):
                    for para in section_content:
                        doc.add_paragraph(para)
            
            # Tables
            for table_data in content.get('tables', []):
                if table_data.get('title'):
                    doc.add_paragraph(table_data['title'], style='Heading 3')
                
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
                
                if headers or rows:
                    table = doc.add_table(
                        rows=1 + len(rows),
                        cols=len(headers) if headers else len(rows[0]) if rows else 1
                    )
                    table.style = 'Table Grid'
                    table.alignment = WD_TABLE_ALIGNMENT.CENTER
                    
                    # Header row
                    if headers:
                        header_cells = table.rows[0].cells
                        for i, header in enumerate(headers):
                            header_cells[i].text = str(header)
                            # Style header
                            for paragraph in header_cells[i].paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                            self._set_cell_shading(header_cells[i], self.style.primary_color)
                    
                    # Data rows
                    for row_idx, row_data in enumerate(rows):
                        row_cells = table.rows[row_idx + 1].cells
                        for col_idx, cell_value in enumerate(row_data):
                            if col_idx < len(row_cells):
                                row_cells[col_idx].text = str(cell_value)
                        
                        # Alternate row shading
                        if row_idx % 2 == 1:
                            for cell in row_cells:
                                self._set_cell_shading(cell, '#f8f9fa')
                    
                    doc.add_paragraph()  # Spacer after table
            
            # Save
            doc.save(output_path)
            
            file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
            elapsed = (time.time() - start_time) * 1000
            
            return ExportResult(
                success=True,
                file_path=output_path,
                format='docx',
                file_size=file_size,
                generation_time_ms=elapsed,
                backend_used='python-docx'
            )
            
        except Exception as e:
            elapsed = (time.time() - start_time) * 1000
            return ExportResult(
                success=False,
                file_path=None,
                format='docx',
                generation_time_ms=elapsed,
                error=str(e)
            )
    
    def _setup_styles(self, doc):
        """Set up custom document styles."""
        from docx.shared import Pt, RGBColor
        
        styles = doc.styles
        
        # Modify heading styles
        for i in range(1, 4):
            style_name = f'Heading {i}'
            if style_name in styles:
                style = styles[style_name]
                style.font.name = self.style.heading_font
                style.font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
    
    def _add_header(self, doc, title: str):
        """Add document header."""
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        section = doc.sections[0]
        header = section.header
        header_para = header.paragraphs[0]
        header_para.text = f"{self.style.company_name} | {title}"
        header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        header_para.style.font.size = Pt(self.style.small_size)
    
    def _add_footer(self, doc):
        """Add document footer with page numbers."""
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        section = doc.sections[0]
        footer = section.footer
        footer_para = footer.paragraphs[0]
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add page number field
        run = footer_para.add_run()
        fldChar1 = OxmlElement('w:fldChar')
        fldChar1.set(qn('w:fldCharType'), 'begin')
        
        instrText = OxmlElement('w:instrText')
        instrText.text = "PAGE"
        
        fldChar2 = OxmlElement('w:fldChar')
        fldChar2.set(qn('w:fldCharType'), 'separate')
        
        fldChar3 = OxmlElement('w:fldChar')
        fldChar3.set(qn('w:fldCharType'), 'end')
        
        run._r.append(fldChar1)
        run._r.append(instrText)
        run._r.append(fldChar2)
        run._r.append(fldChar3)
        
        footer_para.add_run(f" | Generated: {datetime.now().strftime('%Y-%m-%d')}")
    
    def _set_cell_shading(self, cell, color: str):
        """Set cell background color."""
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        shading_elm = OxmlElement('w:shd')
        shading_elm.set(qn('w:fill'), color.replace('#', ''))
        cell._tc.get_or_add_tcPr().append(shading_elm)


# =============================================================================
# PPTX EXPORTER
# =============================================================================

class PPTXExporter:
    """
    Professional PowerPoint exporter with branded templates.
    """
    
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
        
    def export(
        self,
        content: Dict[str, Any],
        output_path: str
    ) -> ExportResult:
        """
        Export content to a professionally styled PPTX.
        
        Args:
            content: Dictionary with presentation content:
                - title: Presentation title
                - subtitle: Optional subtitle
                - author: Author name
                - slides: List of slide dicts with:
                    - layout: 'title', 'section', 'content', 'two_column', 'table', 'chart'
                    - title: Slide title
                    - content: Slide content (varies by layout)
            output_path: Where to save the PPTX
            
        Returns:
            ExportResult with details
        """
        import time
        start_time = time.time()
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            return ExportResult(
                success=False,
                file_path=None,
                format='pptx',
                error="python-pptx not installed. Run: pip install python-pptx"
            )
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # Widescreen 16:9
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._add_title_slide(prs, content)
            
            # Content slides
            for slide_data in content.get('slides', []):
                layout = slide_data.get('layout', 'content')
                
                if layout == 'section':
                    self._add_section_slide(prs, slide_data)
                elif layout == 'two_column':
                    self._add_two_column_slide(prs, slide_data)
                elif layout == 'table':
                    self._add_table_slide(prs, slide_data)
                elif layout == 'metrics':
                    self._add_metrics_slide(prs, slide_data)
                else:
                    self._add_content_slide(prs, slide_data)
            
            # Thank you slide
            self._add_closing_slide(prs, content)
            
            # Save
            prs.save(output_path)
            
            file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
            elapsed = (time.time() - start_time) * 1000
            
            return ExportResult(
                success=True,
                file_path=output_path,
                format='pptx',
                file_size=file_size,
                generation_time_ms=elapsed,
                backend_used='python-pptx'
            )
            
        except Exception as e:
            elapsed = (time.time() - start_time) * 1000
            return ExportResult(
                success=False,
                file_path=None,
                format='pptx',
                generation_time_ms=elapsed,
                error=str(e)
            )
    
    def _add_title_slide(self, prs, content: Dict):
        """Add title slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor.from_string(self.style.primary_color[1:])
        background.line.fill.background()
        
        # Accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(3.5), prs.slide_width, Inches(0.1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor.from_string(self.style.accent_color[1:])
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.get('title', 'Report')
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if content.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.paragraphs[0].text = content['subtitle']
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Date and author
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        info_frame = info_box.text_frame
        author = content.get('metadata', {}).get('author', '')
        date = content.get('metadata', {}).get('date', datetime.now().strftime('%B %Y'))
        info_frame.paragraphs[0].text = f"{author}  |  {date}"
        info_frame.paragraphs[0].font.size = Pt(14)
        info_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_section_slide(self, prs, slide_data: Dict):
        """Add section divider slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Left accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, Inches(0.3), prs.slide_height
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor.from_string(self.style.accent_color[1:])
        accent.line.fill.background()
        
        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get('title', '')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
    
    def _add_content_slide(self, prs, slide_data: Dict):
        """Add standard content slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get('title', '')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
        
        # Title underline
        from pptx.enum.shapes import MSO_SHAPE
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.1), Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(self.style.accent_color[1:])
        line.line.fill.background()
        
        # Content
        content = slide_data.get('content', '')
        if isinstance(content, str):
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.paragraphs[0].text = content
            content_frame.paragraphs[0].font.size = Pt(18)
        elif isinstance(content, list):
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            for i, item in enumerate(content):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.space_after = Pt(8)
    
    def _add_two_column_slide(self, prs, slide_data: Dict):
        """Add two-column slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get('title', '')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_content = slide_data.get('left_content', [])
        for i, item in enumerate(left_content if isinstance(left_content, list) else [left_content]):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = f"• {item}" if isinstance(left_content, list) else str(item)
            p.font.size = Pt(14)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_content = slide_data.get('right_content', [])
        for i, item in enumerate(right_content if isinstance(right_content, list) else [right_content]):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = f"• {item}" if isinstance(right_content, list) else str(item)
            p.font.size = Pt(14)
    
    def _add_table_slide(self, prs, slide_data: Dict):
        """Add slide with table."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get('title', '')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
        
        # Table
        headers = slide_data.get('headers', [])
        rows = slide_data.get('rows', [])
        
        if headers or rows:
            num_cols = len(headers) if headers else len(rows[0]) if rows else 1
            num_rows = (1 if headers else 0) + len(rows)
            
            table = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(0.5), Inches(1.3),
                Inches(12.333), Inches(min(5.5, num_rows * 0.5))
            ).table
            
            # Header row
            if headers:
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = str(header)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(self.style.primary_color[1:])
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(12)
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
            
            # Data rows
            start_row = 1 if headers else 0
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < num_cols:
                        cell = table.cell(start_row + row_idx, col_idx)
                        cell.text = str(cell_value)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(11)
    
    def _add_metrics_slide(self, prs, slide_data: Dict):
        """Add slide with metric cards."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data.get('title', 'Key Metrics')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.style.primary_color[1:])
        
        # Metric cards
        metrics = slide_data.get('metrics', [])
        num_metrics = len(metrics)
        
        if num_metrics > 0:
            card_width = (12.333 - 0.3 * (num_metrics - 1)) / min(num_metrics, 4)
            
            for i, metric in enumerate(metrics[:4]):  # Max 4 per row
                x = 0.5 + i * (card_width + 0.3)
                
                # Card background
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(1.5), Inches(card_width), Inches(2)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor.from_string(self.style.primary_color[1:])
                card.line.fill.background()
                
                # Value
                value_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(1.7), Inches(card_width - 0.2), Inches(1)
                )
                value_frame = value_box.text_frame
                value_frame.paragraphs[0].text = str(metric.get('value', ''))
                value_frame.paragraphs[0].font.size = Pt(36)
                value_frame.paragraphs[0].font.bold = True
                value_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Label
                label_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(2.8), Inches(card_width - 0.2), Inches(0.5)
                )
                label_frame = label_box.text_frame
                label_frame.paragraphs[0].text = metric.get('label', '')
                label_frame.paragraphs[0].font.size = Pt(14)
                label_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
                label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_closing_slide(self, prs, content: Dict):
        """Add closing/thank you slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor.from_string(self.style.primary_color[1:])
        background.line.fill.background()
        
        # Thank you text
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        text_frame = text_box.text_frame
        text_frame.paragraphs[0].text = "Thank You"
        text_frame.paragraphs[0].font.size = Pt(44)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Contact info
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1))
        contact_frame = contact_box.text_frame
        contact_frame.paragraphs[0].text = self.style.company_name
        contact_frame.paragraphs[0].font.size = Pt(18)
        contact_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# =============================================================================
# UNIFIED EXPORT ENGINE
# =============================================================================

class ExportEngine:
    """
    Unified export engine for all formats.
    """
    
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
        self.pdf_exporter = PDFExporter(self.style)
        self.docx_exporter = DOCXExporter(self.style)
        self.pptx_exporter = PPTXExporter(self.style)
        
        # Ensure output directory exists
        self.output_dir = Path(DATA_OUTPUTS) / 'reports'
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def get_available_formats(self) -> Dict[str, bool]:
        """Get available export formats."""
        return {
            'html': True,
            'pdf': len(self.pdf_exporter.available_backends) > 0,
            'docx': True,  # python-docx should be installed
            'pptx': True   # python-pptx should be installed
        }
    
    def export_html(
        self,
        html_content: str,
        filename: Optional[str] = None
    ) -> ExportResult:
        """Export HTML content to file."""
        import time
        start_time = time.time()
        
        if not filename:
            filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        
        output_path = str(self.output_dir / filename)
        
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            file_size = os.path.getsize(output_path)
            elapsed = (time.time() - start_time) * 1000
            
            return ExportResult(
                success=True,
                file_path=output_path,
                format='html',
                file_size=file_size,
                generation_time_ms=elapsed,
                backend_used='native'
            )
        except Exception as e:
            return ExportResult(
                success=False,
                file_path=None,
                format='html',
                error=str(e)
            )
    
    def export_pdf(
        self,
        html_content: str,
        filename: Optional[str] = None,
        preferred_backend: Optional[PDFBackend] = None
    ) -> ExportResult:
        """Export HTML content to PDF."""
        if not filename:
            filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        output_path = str(self.output_dir / filename)
        return self.pdf_exporter.export(html_content, output_path, preferred_backend)
    
    def export_docx(
        self,
        content: Dict[str, Any],
        filename: Optional[str] = None,
        **kwargs
    ) -> ExportResult:
        """Export structured content to DOCX."""
        if not filename:
            filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        output_path = str(self.output_dir / filename)
        return self.docx_exporter.export(content, output_path, **kwargs)
    
    def export_pptx(
        self,
        content: Dict[str, Any],
        filename: Optional[str] = None
    ) -> ExportResult:
        """Export structured content to PPTX."""
        if not filename:
            filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        output_path = str(self.output_dir / filename)
        return self.pptx_exporter.export(content, output_path)
    
    def export_all(
        self,
        html_content: str,
        structured_content: Dict[str, Any],
        base_filename: str
    ) -> Dict[str, ExportResult]:
        """Export to all available formats."""
        results = {}
        
        # HTML
        results['html'] = self.export_html(
            html_content,
            f"{base_filename}.html"
        )
        
        # PDF
        results['pdf'] = self.export_pdf(
            html_content,
            f"{base_filename}.pdf"
        )
        
        # DOCX
        results['docx'] = self.export_docx(
            structured_content,
            f"{base_filename}.docx"
        )
        
        # PPTX (if presentation content provided)
        if 'slides' in structured_content:
            results['pptx'] = self.export_pptx(
                structured_content,
                f"{base_filename}.pptx"
            )
        
        return results


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

_export_engine = None

def get_export_engine(style_config: Optional[StyleConfig] = None) -> ExportEngine:
    """Get or create the export engine singleton."""
    global _export_engine
    if _export_engine is None:
        _export_engine = ExportEngine(style_config)
    return _export_engine


def export_to_pdf(html_content: str, filename: Optional[str] = None) -> ExportResult:
    """Quick export to PDF."""
    return get_export_engine().export_pdf(html_content, filename)


def export_to_docx(content: Dict[str, Any], filename: Optional[str] = None) -> ExportResult:
    """Quick export to DOCX."""
    return get_export_engine().export_docx(content, filename)


def export_to_pptx(content: Dict[str, Any], filename: Optional[str] = None) -> ExportResult:
    """Quick export to PPTX."""
    return get_export_engine().export_pptx(content, filename)


# =============================================================================
# MAIN
# =============================================================================

if __name__ == "__main__":
    print("=" * 60)
    print("TRIALPULSE NEXUS 10X - Export Engine v1.0")
    print("=" * 60)
    
    engine = get_export_engine()
    
    print("\n📋 Available Formats:")
    for fmt, available in engine.get_available_formats().items():
        status = "✅" if available else "❌"
        print(f"   {status} {fmt.upper()}")
    
    print("\n📋 PDF Backends:")
    for backend in engine.pdf_exporter.available_backends:
        print(f"   ✅ {backend.value}")
    
    if not engine.pdf_exporter.available_backends:
        print("   ⚠️  No PDF backends available")
        print("   Install one of: pip install weasyprint / pdfkit / xhtml2pdf")
    
    print("\n✅ Export Engine Ready")